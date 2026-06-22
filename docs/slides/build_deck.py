#!/usr/bin/env python3
"""Generate the Child-Safety-Chatbot presentation as a .pptx.

Run:  .venv/bin/python docs/slides/build_deck.py
Out:  docs/slides/child-safety-chatbot.pptx   (import into Google Slides, fully editable)

The deck argues: a child-safety chatbot is a USE CASE of the compliance-auditor
coding agent — same spine (panel -> threshold -> human-in-the-loop -> hash chain),
new domain.  Slides are credited per teammate as requested.
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")

# ---- theme -----------------------------------------------------------------
BG     = RGBColor(0x0B, 0x12, 0x20)   # near slate-950
PANEL  = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
PANEL2 = RGBColor(0x15, 0x20, 0x33)
BORDER = RGBColor(0x33, 0x41, 0x55)   # slate-700
TXT    = RGBColor(0xE2, 0xE8, 0xF0)   # slate-200
MUTED  = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xEF, 0x44, 0x44)

CAT = {  # matches monitor-web CAT map
    "grooming":          RGBColor(0xA8, 0x55, 0xF7),
    "bullying":          RGBColor(0xF9, 0x73, 0x16),
    "suicidal_ideation": RGBColor(0xEF, 0x44, 0x44),
    "self_harm":         RGBColor(0xBE, 0x12, 0x3C),
    "pii_exposure":      RGBColor(0x3B, 0x82, 0xF6),
    "distress":          RGBColor(0x06, 0xB6, 0xD4),
    "abuse":             RGBColor(0xDB, 0x27, 0x77),  # proposed (not yet scored)
}

FONT = "Arial"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---- helpers ---------------------------------------------------------------
def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, l, t, w, h, fill=PANEL, line=BORDER, line_w=1.0, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, l, t, w, h, runs, size=18, color=TXT, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.0):
    """runs: str, or list of (str, kwargs) paragraphs/lines."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    for i, item in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(item, tuple):
            txt, kw = item
        else:
            txt, kw = item, {}
        # support inline segments via list of (text, kwargs)
        segs = txt if isinstance(txt, list) else [(txt, {})]
        for seg_txt, seg_kw in segs:
            m = {**kw, **seg_kw}   # per-segment kwargs override paragraph kwargs
            r = p.add_run(); r.text = seg_txt
            r.font.name = m.get("font", font)
            r.font.size = Pt(m.get("size", size))
            r.font.bold = m.get("bold", bold)
            r.font.color.rgb = m.get("color", color)
        if "space_after" in kw:
            p.space_after = Pt(kw["space_after"])
    return tb


def chip(s, l, t, label, color, w=None, size=12):
    w = w or (0.16 * len(label) + 0.3)
    c = box(s, l, t, w, 0.34, fill=color, line=None)
    c.text_frame.word_wrap = False
    p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    c.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = p.add_run(); r.text = label
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    return c, w


def arrow(s, l, t, w, h, color=MUTED, shape=MSO_SHAPE.DOWN_ARROW):
    a = box(s, l, t, w, h, fill=color, line=None, rounded=False)
    a.adjustments  # noqa
    a._element.getparent()  # keep ref
    a.fill.solid(); a.fill.fore_color.rgb = color
    a.line.fill.background()
    # replace rounded rect with arrow autoshape
    return a


def autoshape(s, mso, l, t, w, h, fill, line=None):
    shp = s.shapes.add_shape(mso, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def bar(s, l, t, w, frac, color, h=0.16, track=RGBColor(0x33, 0x41, 0x55)):
    box(s, l, t, w, h, fill=track, line=None)
    if frac > 0:
        box(s, l, t, max(0.04, w * frac), h, fill=color, line=None)


def kicker(s, who, section):
    """Top-left teammate + section tag."""
    text(s, 0.55, 0.32, 9, 0.4, [([(who + "  ", {"color": BLUE, "bold": True, "size": 14}),
                                   ("· " + section, {"color": MUTED, "size": 14})], {})], size=14)


def title(s, t, sub=None, top=0.78):
    text(s, 0.55, top, 12.2, 0.8, t, size=32, bold=True, color=WHITE)
    if sub:
        text(s, 0.55, top + 0.62, 12.2, 0.5, sub, size=16, color=MUTED)


# ===========================================================================
# RESEARCH CONTENT  (populated from the research agent; edit then re-run)
# ===========================================================================
WHY_EXAMPLES = [  # synthetic, illustrative — not real victim data
    ("“you seem so mature for 12 — don’t tell your mom we talk, it’s our secret 😊”", "grooming",
     "secrecy request + age + flattery: classic grooming"),
    ("“nobody in this group chat even likes you, just log off forever loser”", "bullying",
     "targeted humiliation and exclusion"),
    ("“honestly I don’t see the point of waking up tomorrow anymore”", "suicidal_ideation",
     "passive suicidal ideation — highest priority"),
    ("“I keep a blade in my drawer for when the bad days get too loud”", "self_harm",
     "reference to means + intent for self-injury"),
    ("“here’s my address, 14 Maple St, and mom’s gone till 6”", "pii_exposure",
     "minor discloses home address + unsupervised window"),
    ("“i can’t stop shaking and crying, everything’s falling apart”", "distress",
     "acute distress without explicit self-harm"),
]

# 10 articles  -> (short title, venue+year, one-line, url)
ARTICLES: list[tuple] = [
    ("PAN-2012 Sexual Predator Identification (overview)", "Inches & Crestani, CLEF/PAN 2012",
     "Defined the canonical grooming-detection benchmark (predator ID + line-level).", ""),
    ("Detecting sexual predators in chats", "Ebrahimi et al., Nat. Lang. Eng. 2016",
     "Behavioral + content features with imbalanced learning on PAN-12.", ""),
    ("Interdisciplinary review of cybergrooming research", "An, Silva, Zhang et al., arXiv 2025",
     "Recent PRISMA survey across social + computational science.", ""),
    ("Early detection of predators via turn-level optimization", "arXiv 2025",
     "Reframes grooming detection as early, turn-level detection.", ""),
    ("Using Machine Learning to Detect Cyberbullying", "Reynolds et al., IEEE ICMLA 2011",
     "Foundational cyberbullying paper; Formspring data via MTurk labels.", ""),
    ("Cyberbullying detection: datasets & approaches (survey)", "arXiv 2024",
     "Recent survey of cyberbullying datasets and ML/DL methods.", ""),
    ("Automated Hate Speech Detection & Offensive Language", "Davidson et al., ICWSM 2017",
     "~25k-tweet hate/offensive/neither dataset; separability problem.", ""),
    ("Crowdsourcing Twitter Abusive Behavior", "Founta et al., ICWSM 2018",
     "~100k-tweet abusive-behavior dataset (hateful/abusive/spam/normal).", ""),
    ("A Survey on Automatic Detection of Hate Speech", "Fortuna & Nunes, ACM CSUR 2018",
     "Widely-cited survey unifying definitions, features, algorithms.", ""),
    ("Suicidal Ideation Detection: an ML review", "Ji et al., IEEE TCSS 2021",
     "Comprehensive ML/DL survey for suicidal-ideation detection.", ""),
]

# datasets -> (name, source/year, size, maps-to, url)
DATASETS: list[tuple] = [
    ("PAN-2012 Sexual Predator ID", "PAN@CLEF 2012", "tens of thousands of chats (<4% predatory)", "grooming", ""),
    ("Formspring.me cyberbullying", "Reynolds 2011", "13,158 posts · 892 bullying", "bullying", ""),
    ("Jigsaw / Wikipedia Toxic Comments", "Jigsaw+Google 2018", "~159,571 comments · 6 toxicity labels", "bullying / toxicity", ""),
    ("OLID (OffensEval, SemEval-2019 T6)", "Zampieri 2019", "~14,100 tweets · 3-level", "bullying / abusive", ""),
    ("Davidson Hate Speech & Offensive", "Davidson 2017", "~24,800 tweets", "bullying / abusive", ""),
    ("CLPsych-2016 (ReachOut forum)", "Milne 2016", "~65k posts · triage green→crisis", "suicidal / distress", ""),
    ("Reddit SuicideWatch", "Kaggle 2008–21", "232,074 posts (balanced)", "suicidal / self-harm", ""),
]
# why-it-matters facts -> (stat, source/year)
WHY_FACTS: list[tuple] = [
    ("36.2M+ CyberTipline reports in 2023 (>100M files)", "NCMEC 2023 CyberTipline Report"),
    ("UK Online Safety Act 2023 — duty of care; fines up to 10% of global turnover", "UK legislation, in force 2023"),
    ("EU Child Sexual Abuse Regulation — Council position reached Nov 2025", "European Commission / Council"),
    ("720,000+ suicide deaths/yr; 3rd leading cause of death ages 15–29", "WHO fact sheet, 2025 (2021 data)"),
    ("46% of US teens have experienced cyberbullying", "Pew Research Center, 2022"),
]


# ===========================================================================
# SLIDES
# ===========================================================================
def s_title():
    s = slide()
    box(s, 0.55, 0.55, 0.9, 0.9, fill=BLUE, line=None)
    text(s, 0.55, 0.55, 0.9, 0.9, "SG", size=30, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.7, 0.62, 10, 0.8, "SSGCheck", size=20, bold=True, color=WHITE)
    text(s, 1.7, 1.02, 10, 0.5, "non-repudiable oversight for AI", size=14, color=MUTED)

    text(s, 0.55, 2.5, 12.2, 1.6,
         [("A Child-Safety Chatbot", {"size": 46, "bold": True, "color": WHITE}),
          ("as a use case of a compliance-auditor coding agent", {"size": 30, "bold": True, "color": BLUE})],
         line_spacing=1.05)
    text(s, 0.55, 4.5, 12, 0.6,
         "Same spine — panel of judges → threshold → human-in-the-loop → tamper-evident hash chain — pointed at children’s conversations instead of code.",
         size=16, color=MUTED)
    # team
    text(s, 0.55, 6.35, 12, 0.5,
         [([("Team   ", {"color": MUTED, "size": 14}),
            ("Elaine · Grace · Wilson · Ethan · Jeremy", {"color": TXT, "size": 14, "bold": True})], {})])


def s_architecture():
    """The hero: LEFT coding-agent flow, RIGHT child-safety detection, MCP -> child case."""
    s = slide()
    kicker(s, "Jeremy", "System Architecture")
    title(s, "Child-safety chatbot = a use case of the coding agent",
          "Left: the governance spine.  Right: the same spine instantiated for conversations.")

    # ---- LEFT column: coding agent flow ----
    lx, lw = 0.55, 5.7
    text(s, lx, 1.9, lw, 0.4, "THE CODING AGENT (governance spine)", size=13, bold=True, color=BLUE)
    steps = [
        ("Developer request → Coding agent + AGENTS.md", "the “gate”", PANEL),
        ("GATE: don’t code yet · ≥2 options + tradeoffs ·\ncite EU AI Act verbatim · WAIT for human", "prompt-enforced", PANEL),
        ("MCP tool  log_decision(…)", "deterministic Python", RGBColor(0x1d, 0x4e, 0xd8)),
        ("audit-trail/decisions.jsonl  —  SHA-256 hash chain", "append-only, verifiable", PANEL),
    ]
    y = 2.35
    heights = [0.7, 0.95, 0.62, 0.7]
    for idx, ((label, tag, fill), hh) in enumerate(zip(steps, heights)):
        b = box(s, lx, y, lw, hh, fill=fill, line=BORDER)
        b.text_frame.word_wrap = True
        b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = label; r.font.size = Pt(13.5); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = FONT
        p2 = b.text_frame.add_paragraph(); rr = p2.add_run(); rr.text = tag
        rr.font.size = Pt(10); rr.font.color.rgb = MUTED; rr.font.name = FONT
        y += hh + 0.28
        if idx < len(steps) - 1:
            autoshape(s, MSO_SHAPE.DOWN_ARROW, lx + lw/2 - 0.12, y - 0.27, 0.24, 0.2, MUTED)
    text(s, lx, y - 0.02, lw, 0.5,
         "└▶ verify_audit_trail()    └▶ generate_compliance_report()",
         size=11, color=GREEN)

    # ---- connector: MCP -> child-safety case ----
    # arrow from log_decision box (left) across to right column
    autoshape(s, MSO_SHAPE.RIGHT_ARROW, lx + lw + 0.02, 3.95, 1.05, 0.5, BLUE)
    text(s, lx + lw - 0.1, 4.5, 1.4, 0.4, "same log,\nnew domain", size=9, color=BLUE, align=PP_ALIGN.CENTER)

    # ---- RIGHT column: child-safety detection ----
    rx, rw = 7.5, 5.3
    text(s, rx, 1.9, rw, 0.4, "THE CHILD-SAFETY CHATBOT (detection)", size=13, bold=True, color=RED)
    det = [
        ("Child / Adult conversation turns", "input", PANEL),
        ("3-LLM risk panel  (DeepSeek · OpenAI · Gemini)\nscore 6 categories 0–1", "ensemble", PANEL),
        ("aggregate ≥ threshold ?  →  ALARM", "policy", RGBColor(0x7f, 0x1d, 0x1d)),
        ("ESCALATE to human monitor  →  log_decision(domain=“child-safety”)", "same hash chain", RGBColor(0x1d, 0x4e, 0xd8)),
    ]
    y = 2.35
    hh2 = [0.62, 0.85, 0.6, 0.8]
    for idx, ((label, tag, fill), h) in enumerate(zip(det, hh2)):
        b = box(s, rx, y, rw, h, fill=fill, line=BORDER)
        b.text_frame.word_wrap = True; b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = b.text_frame.paragraphs[0]
        r = p.add_run(); r.text = label; r.font.size = Pt(13.5); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = FONT
        p2 = b.text_frame.add_paragraph(); rr = p2.add_run(); rr.text = tag
        rr.font.size = Pt(10); rr.font.color.rgb = MUTED; rr.font.name = FONT
        y += h + 0.26
        if idx < len(det) - 1:
            autoshape(s, MSO_SHAPE.DOWN_ARROW, rx + rw/2 - 0.12, y - 0.25, 0.24, 0.18, MUTED)
    # category chips row (wrap before overflowing the right edge)
    cy = y + 0.05
    cx = rx
    chips = [("grooming", "grooming"), ("bullying", "bullying"), ("suicidal", "suicidal_ideation"),
             ("self-harm", "self_harm"), ("PII", "pii_exposure"), ("distress", "distress")]
    for disp, key in chips:
        w = 0.115 * len(disp) + 0.34
        if cx + w > rx + rw + 0.05:
            cx = rx; cy += 0.46
        chip(s, cx, cy, disp, CAT[key], w=w, size=10)
        cx += w + 0.12


def s_why():
    s = slide()
    kicker(s, "Elaine", "Why it matters")
    title(s, "Why a child-safety chatbot matters",
          "Conversations with minors carry risks that are invisible to keyword filters — we label them.")
    # left: labeled examples
    lx = 0.55
    text(s, lx, 1.95, 6.4, 0.4, "UNSAFE MESSAGES, LABELED", size=13, bold=True, color=BLUE)
    short = {"grooming": "grooming", "bullying": "bullying", "suicidal_ideation": "suicidal",
             "self_harm": "self-harm", "pii_exposure": "PII", "distress": "distress"}
    y = 2.3
    for txt_, label, why in WHY_EXAMPLES:
        box(s, lx, y, 6.45, 0.78, fill=PANEL, line=BORDER)
        text(s, lx + 0.18, y + 0.07, 6.1, 0.4, txt_, size=11.5, color=TXT)
        w = 0.12 * len(short[label]) + 0.34
        chip(s, lx + 0.18, y + 0.43, short[label], CAT.get(label, MUTED), w=w, size=10)
        text(s, lx + 0.18 + w + 0.16, y + 0.46, 6.45 - (w + 0.6), 0.3, why, size=9.5, color=MUTED)
        y += 0.86
    # right: why-now facts
    rx = 7.4
    text(s, rx, 1.95, 5.3, 0.4, "WHY NOW — SCALE & REGULATION", size=13, bold=True, color=RED)
    yy = 2.45
    if WHY_FACTS:
        for stat, src in WHY_FACTS:
            box(s, rx, yy, 0.12, 0.12, fill=RED, line=None)
            text(s, rx + 0.28, yy - 0.06, 5.0, 0.7, [(stat, {"size": 13, "color": TXT, "bold": True}),
                                                     (src, {"size": 10, "color": MUTED})])
            yy += 0.95
    else:
        text(s, rx, yy, 5.0, 2, "[ research agent: regulatory + scale facts ]", size=12, color=MUTED)


def s_dataset():
    s = slide()
    kicker(s, "Elaine", "Research landscape & datasets")
    title(s, "A trending research topic — with public benchmarks",
          "Detection of grooming, cyberbullying and suicidal ideation is an active NLP area with established datasets.")
    # left: articles
    lx = 0.55
    text(s, lx, 1.95, 6.3, 0.4, "SELECTED LITERATURE", size=13, bold=True, color=BLUE)
    y = 2.4
    if ARTICLES:
        for (t, venue, one, url) in ARTICLES[:8]:
            text(s, lx, y, 6.3, 0.6,
                 [([(t + "  ", {"size": 11.5, "color": TXT, "bold": True}),
                    ("(" + venue + ")", {"size": 10, "color": MUTED})], {}),
                  (one, {"size": 10, "color": MUTED})], line_spacing=0.95)
            y += 0.6
    else:
        text(s, lx, y, 6.3, 3, "[ research agent: 10 articles ]", size=12, color=MUTED)
    # right: datasets
    rx = 7.2
    text(s, rx, 1.95, 5.6, 0.4, "BENCHMARK DATASETS", size=13, bold=True, color=RED)
    y = 2.4
    if DATASETS:
        for (name, src, size_, maps, url) in DATASETS[:7]:
            b = box(s, rx, y, 5.6, 0.62, fill=PANEL, line=BORDER)
            text(s, rx + 0.16, y + 0.06, 4.0, 0.4,
                 [([(name + "  ", {"size": 11.5, "bold": True, "color": TXT}),
                    (src, {"size": 9.5, "color": MUTED})], {}),
                  (size_ + "  ·  " + maps, {"size": 9.5, "color": MUTED})], line_spacing=0.92)
            y += 0.72
    else:
        text(s, rx, y, 5.6, 3, "[ research agent: datasets ]", size=12, color=MUTED)


def s_risks():
    s = slide()
    kicker(s, "Grace", "How risks are defined")
    title(s, "Six categories, anchored to UN documents",
          "Each category maps to a Convention on the Rights of the Child article + the EU AI Act.")
    rows = [
        ("grooming", "Grooming / sexual solicitation", "CRC Art. 34 (sexual exploitation & abuse), Art. 19"),
        ("bullying", "Bullying / peer aggression", "CRC Art. 19 (protection from violence), Art. 3"),
        ("suicidal_ideation", "Suicidal ideation", "CRC Art. 6 (life, survival), Art. 3 (best interests)"),
        ("self_harm", "Self-harm", "CRC Art. 6, Art. 24 (health)"),
        ("pii_exposure", "PII exposure", "CRC Art. 16 (privacy)"),
        ("distress", "Emotional distress / help-seeking", "CRC Art. 3, Art. 12 (right to be heard)"),
    ]
    # header
    y = 2.02
    text(s, 0.7, y, 3.2, 0.3, "CATEGORY", size=11, bold=True, color=MUTED)
    text(s, 4.0, y, 3.6, 0.3, "WHAT IT CAPTURES", size=11, bold=True, color=MUTED)
    text(s, 7.7, y, 5.2, 0.3, "LEGAL ANCHOR (verbatim, not LLM-recalled)", size=11, bold=True, color=MUTED)
    y += 0.36
    for k, what, law in rows:
        box(s, 0.55, y, 12.25, 0.56, fill=PANEL2, line=BORDER)
        box(s, 0.7, y + 0.15, 0.26, 0.26, fill=CAT[k], line=None)
        text(s, 1.1, y + 0.11, 2.8, 0.4, k.replace("_", " "), size=13, bold=True, color=TXT)
        text(s, 4.0, y + 0.11, 3.6, 0.4, what, size=12, color=TXT)
        text(s, 7.7, y + 0.11, 5.0, 0.5, law, size=11, color=MUTED)
        y += 0.62
    # overall risk definition
    box(s, 0.55, y + 0.05, 12.25, 0.92, fill=RGBColor(0x12, 0x1c, 0x2e), line=BLUE)
    text(s, 0.8, y + 0.15, 12, 0.8,
         [([("Overall risk   ", {"size": 13, "bold": True, "color": BLUE}),
            ("each of 3 models scores every category 0–1;  ", {"size": 12, "color": TXT}),
            ("by_measure", {"size": 12, "bold": True, "color": WHITE}),
            (" = worst score across the panel,  ", {"size": 12, "color": TXT}),
            ("aggregate", {"size": 12, "bold": True, "color": WHITE}),
            (" = mean of model overalls.", {"size": 12, "color": TXT})], {}),
          ("ALARM when aggregate ≥ threshold (0.6) → escalate to a human + append to the hash chain.",
           {"size": 12, "color": MUTED})])


def s_prompts():
    s = slide()
    kicker(s, "Wilson", "System prompt development")
    title(s, "Three system prompts",
          "One governance prompt + one scoring prompt + the structured-output contract that makes it labelable.")
    cards = [
        ("1 · The gate  (AGENTS.md)", BLUE,
         ["Injected into every coding-agent call.",
          "Defines WHEN the gate fires (privacy / security /",
          "fairness / child-safety) and WHAT to do: present",
          "≥2 options + tradeoffs, cite CRC/AI Act verbatim,",
          "WAIT for the human, then log_decision()."]),
        ("2 · The risk scorer  (score_one)", RED,
         ["“You are a strict AI-risk auditor. Score the given",
          "{kind} for risk on each measure 0.0–1.0.”",
          "Parameterised by kind + measures — the SAME",
          "template scores code and conversations.",
          "Run independently by all 3 models."]),
        ("3 · The output contract  (JSON)", GREEN,
         ["Reply with ONLY JSON: {scores:{<measure>:float},",
          "overall:float, reason:‘…’}.",
          "Forces a machine-labelable result → per-message",
          "labels + the distribution visualization.",
          "Floor 0.15 → ‘none’ (no salient risk)."]),
    ]
    x = 0.55
    for head, accent, lines in cards:
        b = box(s, x, 2.2, 4.0, 4.2, fill=PANEL, line=BORDER)
        box(s, x, 2.2, 4.0, 0.12, fill=accent, line=None, rounded=False)
        text(s, x + 0.25, 2.45, 3.6, 0.5, head, size=15, bold=True, color=WHITE)
        text(s, x + 0.25, 3.1, 3.55, 3.2,
             [(ln, {"size": 11.5, "color": TXT if not ln.startswith("“") else MUTED}) for ln in lines],
             line_spacing=1.18)
        x += 4.2
    text(s, 0.55, 6.6, 12, 0.5,
         "Child-safety adds no new prompt — it is a domain in prompt 1 and a measure-set in prompt 2.  That is what makes it a use case, not a separate system.",
         size=12, color=MUTED)


def s_demo():
    s = slide()
    kicker(s, "Ethan", "System demo")
    title(s, "Live demo — per-message labels + risk distribution")
    # screenshot
    pic = os.path.join(IMG, "demo-overview.png")
    if os.path.exists(pic):
        s.shapes.add_picture(pic, Inches(0.55), Inches(1.85), width=Inches(8.4))
    # callout boxes on the right
    cx = 9.25
    callouts = [
        ("Per-message chip", "every bubble is labeled with its dominant risk category + score", CAT["grooming"]),
        ("3-model panel", "DeepSeek / OpenAI / Gemini score independently → aggregate 0.71 → HIGH", RED),
        ("Risk distribution", "stacked share bar + per-category severity & message counts", BLUE),
        ("Escalation + audit", "alarm → escalate to human → SHA-256 record appended", GREEN),
    ]
    y = 1.95
    for head, body, accent in callouts:
        b = box(s, cx, y, 3.55, 1.15, fill=PANEL, line=BORDER)
        box(s, cx, y + 0.0, 0.1, 1.15, fill=accent, line=None, rounded=False)
        text(s, cx + 0.25, y + 0.12, 3.2, 0.4, head, size=13, bold=True, color=WHITE)
        text(s, cx + 0.25, y + 0.5, 3.2, 0.6, body, size=10.5, color=MUTED, line_spacing=0.98)
        y += 1.28


def s_close():
    s = slide()
    text(s, 0.55, 2.4, 12.2, 1.4,
         [("Child Safety chatbot", {"size": 40, "bold": True, "color": WHITE}),
          ("is a use case of the compliance-auditor coding agent.", {"size": 28, "bold": True, "color": BLUE})],
         line_spacing=1.1)
    pts = [
        "Same spine: panel of judges → threshold → human-in-the-loop → tamper-evident hash chain.",
        "New domain: 6 risk categories, anchored to UN CRC + EU AI Act, scored by the same prompt.",
        "Result: a labeled, visualized, non-repudiable record of every escalation — reusable, not bespoke.",
    ]
    y = 4.4
    for p in pts:
        box(s, 0.6, y + 0.06, 0.12, 0.12, fill=GREEN, line=None)
        text(s, 0.9, y - 0.04, 11.5, 0.5, p, size=15, color=TXT)
        y += 0.6


# textbox helper doesn't accept italic_ok; patch text() signature gracefully
def _patch():
    pass


for fn in (s_title, s_architecture, s_why, s_dataset, s_risks, s_prompts, s_demo, s_close):
    fn()

out = os.path.join(HERE, "child-safety-chatbot.pptx")
prs.save(out)
print("wrote", out, "slides:", len(prs.slides._sldIdLst))
