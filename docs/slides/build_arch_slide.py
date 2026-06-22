#!/usr/bin/env python3
"""One slide: System Architecture — the coding-agent flow, MCP -> child-safety use case.
Run: .venv/bin/python docs/slides/build_arch_slide.py  ->  docs/slides/system-architecture.pptx
"""
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BG     = RGBColor(0x0B, 0x12, 0x20)
PANEL  = RGBColor(0x1E, 0x29, 0x3B)
BORDER = RGBColor(0x33, 0x41, 0x55)
TXT    = RGBColor(0xE2, 0xE8, 0xF0)
MUTED  = RGBColor(0x94, 0xA3, 0xB8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
RED    = RGBColor(0xEF, 0x44, 0x44)
FONT = "Arial"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background(); bg.shadow.inherit = False


def box(l, t, w, h, fill=PANEL, line=BORDER, rounded=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def label(shp, main, sub=None, size=14):
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = main; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = FONT
    if sub:
        p2 = tf.add_paragraph(); rr = p2.add_run(); rr.text = sub
        rr.font.size = Pt(10.5); rr.font.color.rgb = MUTED; rr.font.name = FONT


def txt(l, t, w, h, runs, size=14, color=TXT, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, ls=1.0):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    if isinstance(runs, str): runs = [(runs, {})]
    for i, (tt, kw) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls
        r = p.add_run(); r.text = tt
        r.font.size = Pt(kw.get("size", size)); r.font.bold = kw.get("bold", bold)
        r.font.color.rgb = kw.get("color", color); r.font.name = FONT
    return tb


def down(cx, top):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(cx - 0.11), Inches(top), Inches(0.22), Inches(0.2))
    a.fill.solid(); a.fill.fore_color.rgb = MUTED; a.line.fill.background(); a.shadow.inherit = False


# ---- title ----
txt(0.55, 0.42, 12, 0.8, "System Architecture", size=34, bold=True, color=WHITE)

# ---- left flow ----
lx, lw = 0.55, 8.0
cx = lx + lw / 2
steps = [
    ("developer request", None, 0.5),
    ("Coding agent (OpenCode / any MCP host)  +  AGENTS.md  — the “gate”",
     "1. sensitive signal?  (intent / filename / schema keywords)", 0.78),
    ("GATE (prompt-enforced):  don’t code yet · present ≥2 options + tradeoffs ·\ncite EU AI Act ref verbatim · WAIT for human",
     "2. human chooses", 0.92),
    ("MCP tool  log_decision(…)   →   compliance-auditor (deterministic Python)",
     "3. append-only", 0.66),
    ("audit-trail/decisions.jsonl   —   SHA-256 hash chain\nseq N: { …, prev_hash, hash = sha256(canonical(record)) }",
     None, 0.82),
]
fills = [PANEL, PANEL, PANEL, RGBColor(0x1D, 0x4E, 0xD8), PANEL]
y = 1.5
mcp_box_mid = None
for i, ((main, sub, h), fill) in enumerate(zip(steps, fills)):
    b = box(lx, y, lw, h, fill=fill)
    label(b, main, sub)
    if i == 3:
        mcp_box_mid = y + h / 2          # remember MCP row centre for the side arrow
        mcp_right = lx + lw
    y += h
    if i < len(steps) - 1:
        down(cx, y + 0.03); y += 0.28
# verify / report
txt(lx, y + 0.06, lw, 0.4,
    [("├▶ verify_audit_trail()      └▶ generate_compliance_report()", {"size": 11.5, "color": GREEN})])

# ---- arrow from MCP to the right ----
ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(mcp_right + 0.05), Inches(mcp_box_mid - 0.28),
                        Inches(1.15), Inches(0.56))
ar.fill.solid(); ar.fill.fore_color.rgb = BLUE; ar.line.fill.background(); ar.shadow.inherit = False
txt(mcp_right + 0.0, mcp_box_mid + 0.32, 1.3, 0.3, [("use case", {"size": 10, "color": BLUE, "bold": True})],
    align=PP_ALIGN.CENTER)

# ---- right: use-case text box ----
rx, rw = mcp_right + 1.35, 3.1
rb = box(rx, mcp_box_mid - 1.35, rw, 2.7, fill=RGBColor(0x12, 0x1C, 0x2E), line=BLUE)
tf = rb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.18)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "Child-safety chatbot"; r.font.size = Pt(22); r.font.bold = True
r.font.color.rgb = WHITE; r.font.name = FONT
for line, c, sz in [("a use case of the above", BLUE, 14),
                    ("grooming · bullying · abuse ·", MUTED, 12),
                    ("self-harm · suicidal · distress", MUTED, 12),
                    ("→ same gate, same hash chain", TXT, 12)]:
    pp = tf.add_paragraph(); pp.alignment = PP_ALIGN.LEFT
    rr = pp.add_run(); rr.text = line; rr.font.size = Pt(sz); rr.font.bold = (c is BLUE)
    rr.font.color.rgb = c; rr.font.name = FONT

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "system-architecture.pptx")
prs.save(out); print("wrote", out)
